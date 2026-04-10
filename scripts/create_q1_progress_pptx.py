from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORT_DIR = ROOT / "松山湖" / "汇报"
TEMPLATE_PPTX = REPORT_DIR / "20251104 1131-课题2-浙大部分-发送版本.pptx"
OUTPUT_PPTX = REPORT_DIR / "20260412_q1_project_discussion_slides.pptx"


COLORS = {
    "white": "FFFFFF",
    "black": "000000",
    "blue": "0B4AA2",
    "arrow_blue": "5B9BD5",
    "gray": "7A7A7A",
    "gray_light": "F2F2F2",
    "orange": "FFC000",
    "pink": "E6B8E7",
    "blue_dash": "6FA8FF",
    "purple": "7030A0",
    "beige": "FFF2CC",
}

FONT_HEAD = "SimHei"
FONT_BODY = "Microsoft YaHei"
FONT_SERIF = "SimSun"
SLIDE_W = None
SLIDE_H = None


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_fill(shape, color=None):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color, width=1.2, dash=None):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash is not None:
        shape.line.dash_style = dash


def clear_text_frame(text_frame, margin=0.04):
    text_frame.clear()
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)
    text_frame.word_wrap = True


def set_shape_text(shape, text, font_size, color="black", bold=False, align=PP_ALIGN.LEFT,
                   font_name=FONT_BODY, margin=0.04, valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    clear_text_frame(tf, margin=margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return shape


def add_textbox(slide, x, y, w, h, text, font_size, color="black", bold=False,
                align=PP_ALIGN.LEFT, font_name=FONT_BODY, margin=0.02,
                valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_shape_text(shape, text, font_size, color, bold, align, font_name, margin, valign)


def add_box(slide, x, y, w, h, fill=None, line="gray", width=1.2, dash=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, COLORS[fill] if fill else None)
    set_line(shape, COLORS[line], width, dash)
    return shape


def add_label_tag(slide, x, y, w, h, text, fill, text_color="white"):
    shape = add_box(slide, x, y, w, h, fill=fill, line=fill, width=1.0)
    set_shape_text(shape, text, 16, text_color, True, PP_ALIGN.CENTER, FONT_HEAD, margin=0.01)
    return shape


def add_section_header(slide, page_no):
    add_textbox(slide, 0.55, 0.14, 0.38, 0.45, "1", 34, "blue", True, font_name=FONT_HEAD, margin=0.0)
    add_textbox(slide, 0.95, 0.14, 2.2, 0.45, "研究进展", 31, "blue", True, font_name=FONT_HEAD, margin=0.0)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.86), SLIDE_W, Inches(0.035))
    set_fill(line, COLORS["blue"])
    line.line.fill.background()
    add_textbox(slide, 12.82, 7.03, 0.28, 0.22, str(page_no), 18, "blue", True, align=PP_ALIGN.RIGHT,
                font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)


def add_bullet_heading(slide, y, text, font_size):
    add_textbox(slide, 0.6, y, 12.1, 0.38, f"◆ {text}", font_size, "black", True,
                font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)


def add_arrow(slide, x, y, w, h, direction="right"):
    shape_type = MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW if direction == "right" else MSO_AUTO_SHAPE_TYPE.DOWN_ARROW
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, COLORS["arrow_blue"])
    set_line(shape, COLORS["blue"], 1.0)
    return shape


def add_square_bullets(slide, x, y, w, items, font_size=14.5, gap=0.44, square_color="blue"):
    current_y = y
    for item in items:
        add_textbox(slide, x, current_y, 0.16, 0.22, "▪", font_size + 1, square_color, True,
                    font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)
        add_textbox(slide, x + 0.2, current_y, w - 0.2, 0.28, item, font_size, "black", False,
                    font_name=FONT_BODY, margin=0.0, valign=MSO_ANCHOR.TOP)
        current_y += gap


def add_table_cell(slide, x, y, w, h, text, fill=None, font_size=12, bold=False,
                   align=PP_ALIGN.LEFT, font_name=FONT_BODY):
    cell = add_box(slide, x, y, w, h, fill=fill, line="gray", width=0.8)
    set_shape_text(cell, text, font_size, "black", bold, align, font_name, margin=0.03)
    return cell


def add_module_table(slide, x, y):
    headers = ["单元类型", "典型设备", "能量方向"]
    rows = [
        ("动力发电单元", "CHP、燃气轮机", "天然气→电+热"),
        ("余热利用单元", "吸收式制冷机、余热锅炉", "余热→冷/热"),
        ("风力发电单元", "风电机组", "风能→电"),
        ("光伏发电单元", "光伏机组", "光能→电"),
        ("燃料制备单元", "电解槽、储氢罐", "电→氢"),
        ("电热转化单元", "热泵、电制冷机", "电→热/冷"),
        ("化学储能单元", "蓄电池、蓄热/蓄冷", "电/热/冷储存"),
        ("可再生供冷热单元", "地源热泵、光热集热器", "地热/光热→热/冷"),
    ]
    col_widths = [2.2, 2.45, 2.2]
    row_h = 0.4
    current_x = x
    for idx, header in enumerate(headers):
        add_table_cell(slide, current_x, y, col_widths[idx], row_h, header, fill="gray_light",
                       font_size=13, bold=True, align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
        current_x += col_widths[idx]
    current_y = y + row_h
    for row in rows:
        current_x = x
        for idx, value in enumerate(row):
            align = PP_ALIGN.CENTER if idx == 2 else PP_ALIGN.LEFT
            add_table_cell(slide, current_x, current_y, col_widths[idx], row_h, value,
                           font_size=12, align=align)
            current_x += col_widths[idx]
        current_y += row_h


def add_flow_step(slide, x, y, w, h, title, desc):
    box = add_box(slide, x, y, w, h, fill=None, line="gray", width=1.0)
    set_shape_text(box, f"{title}\n{desc}", 12.5, "black", True, PP_ALIGN.CENTER, FONT_BODY, margin=0.03)
    return box


def add_problem_step(slide, x, y, w, h, text):
    box = add_box(slide, x, y, w, h, fill=None, line="gray", width=1.0, dash=MSO_LINE_DASH_STYLE.DASH)
    set_shape_text(box, text, 16, "black", False, PP_ALIGN.CENTER, FONT_BODY, margin=0.04)
    return box


def add_layer_row(slide, x, y, w, h, label, title, desc):
    outer = add_box(slide, x, y, w, h, fill=None, line="gray", width=1.0)
    label_w = 1.02
    label_box = add_box(slide, x, y, label_w, h, fill="gray_light", line="gray", width=1.0)
    set_shape_text(label_box, label, 18, "black", False, PP_ALIGN.CENTER, FONT_HEAD, margin=0.02)

    title_box = add_box(slide, x + 1.12, y + 0.12, 1.7, h - 0.24, fill=None, line="gray", width=0.9,
                        dash=MSO_LINE_DASH_STYLE.DASH)
    set_shape_text(title_box, title, 15, "black", False, PP_ALIGN.CENTER, FONT_BODY, margin=0.02)

    add_arrow(slide, x + 2.9, y + 0.36, 0.38, 0.22, "right")

    desc_box = add_box(slide, x + 3.42, y + 0.06, w - 3.62, h - 0.12, fill=None, line="gray", width=0.9,
                       dash=MSO_LINE_DASH_STYLE.DASH)
    set_shape_text(desc_box, desc, 15, "black", False, PP_ALIGN.CENTER, FONT_BODY, margin=0.04)
    return outer


def add_info_section(slide, x, y, title, desc):
    add_textbox(slide, x, y, 2.3, 0.26, f"□ {title}", 18, "blue", True,
                font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)
    add_textbox(slide, x + 0.08, y + 0.34, 2.85, 0.36, desc, 14.2, "black", False,
                font_name=FONT_BODY, margin=0.0, valign=MSO_ANCHOR.TOP)


def build_slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, 2)

    add_bullet_heading(slide, 1.08, "已完成多能转化单元模块库的梳理与统一接口封装", 27)
    add_bullet_heading(slide, 1.58, "当前模块库已可支撑规划模型快速拼装与设备扩展", 25)

    add_box(slide, 0.72, 2.18, 7.2, 4.9, fill=None, line="gray", width=1.6)
    add_label_tag(slide, 1.18, 2.05, 1.75, 0.24, "模块库清单", "purple")
    add_module_table(slide, 0.94, 2.48)

    interface_box = add_box(slide, 0.96, 6.05, 6.72, 0.65, fill="gray_light", line="gray", width=0.9)
    set_shape_text(interface_box, "统一接口：参数类 / create_component() / get_io_description()",
                   13.5, "black", False, PP_ALIGN.CENTER, FONT_BODY, margin=0.03)

    add_box(slide, 8.18, 2.18, 4.45, 2.32, fill=None, line="gray", width=1.6)
    add_label_tag(slide, 8.58, 2.05, 2.05, 0.24, "接入规划模型", "beige", text_color="black")
    add_flow_step(slide, 8.45, 2.92, 1.14, 0.96, "资源/负荷数据", "风速、辐照、\n气象、园区负荷")
    add_arrow(slide, 9.66, 3.26, 0.3, 0.2, "right")
    add_flow_step(slide, 10.01, 2.92, 1.18, 0.96, "模块参数与IO", "统一参数类、\n组件生成、IO说明")
    add_arrow(slide, 11.27, 3.26, 0.3, 0.2, "right")
    add_flow_step(slide, 11.62, 2.92, 0.86, 0.96, "规划/调度模型", "接入 OEMOF，\n支撑容量规划")
    add_textbox(slide, 8.45, 4.08, 4.0, 0.22, "形成“数据—模块—模型”统一链路", 13.2,
                "black", False, align=PP_ALIGN.CENTER, font_name=FONT_BODY, margin=0.0,
                valign=MSO_ANCHOR.TOP)

    add_box(slide, 8.18, 4.82, 4.45, 2.26, fill=None, line="blue_dash", width=1.5,
            dash=MSO_LINE_DASH_STYLE.DASH)
    add_textbox(slide, 8.48, 5.08, 2.2, 0.28, "后续补充数据", 20, "blue", True,
                font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)
    add_square_bullets(
        slide,
        8.55,
        5.46,
        3.78,
        [
            "设备额定参数、效率曲线、启停边界和典型工况。",
            "园区实测负荷、气象数据和更完整的设备清单。",
            "统一命名与接口字段，便于模块复用与后续维护。",
        ],
        font_size=14.2,
        gap=0.46,
    )
    return slide


def build_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, 3)

    add_bullet_heading(slide, 1.08, "源荷匹配研究已形成“规划—设备—运行”三层递进主线", 27)
    add_bullet_heading(slide, 1.58, "以㶲加权匹配度为起点，逐步扩展到卡诺电池规划与运行层验证", 24.5)

    add_box(slide, 0.2, 2.12, 2.78, 4.95, fill=None, line="orange", width=1.6,
            dash=MSO_LINE_DASH_STYLE.DASH)
    add_label_tag(slide, 0.52, 1.99, 1.25, 0.24, "研究问题", "beige", text_color="black")
    problem_top = add_box(slide, 0.45, 2.48, 2.28, 0.84, fill=None, line="gray", width=1.0)
    set_shape_text(problem_top, "电/热/冷等权处理\n难体现能质差异", 16, "black", False,
                   PP_ALIGN.CENTER, FONT_BODY, margin=0.03)
    add_arrow(slide, 1.38, 3.42, 0.32, 0.3, "down")
    add_problem_step(slide, 0.45, 3.78, 2.28, 0.92, "新型电热耦合储能\n配置价值待评估")
    add_arrow(slide, 1.38, 4.82, 0.32, 0.3, "down")
    add_problem_step(slide, 0.45, 5.18, 2.28, 0.92, "运行阶段存在预测误差\n需验证方案鲁棒性")

    add_box(slide, 3.28, 2.12, 6.0, 4.95, fill=None, line="pink", width=1.6,
            dash=MSO_LINE_DASH_STYLE.DASH)
    title_bar = add_box(slide, 3.48, 2.38, 5.6, 0.34, fill=None, line="gray", width=1.0)
    set_shape_text(title_bar, "源荷匹配研究逻辑主线", 20, "black", False,
                   PP_ALIGN.CENTER, FONT_SERIF, margin=0.02)

    add_layer_row(
        slide,
        3.48,
        2.86,
        5.6,
        1.0,
        "规划层",
        "核心问题\n与方法",
        "引入卡诺㶲系数，构建 EQD 匹配度；\n与年化成本形成双目标容量优化",
    )
    add_arrow(slide, 5.98, 3.9, 0.34, 0.28, "down")
    add_layer_row(
        slide,
        3.48,
        4.08,
        5.6,
        1.0,
        "设备层",
        "模型扩展\n与对比",
        "将卡诺电池纳入 CCHP 模型；\n对比有/无卡诺电池及 EQD/Std 配置差异",
    )
    add_arrow(slide, 5.98, 5.12, 0.34, 0.28, "down")
    add_layer_row(
        slide,
        3.48,
        5.3,
        5.6,
        1.0,
        "运行层",
        "调度验证\n与鲁棒性",
        "结合负荷预测与日前调度；\n对比不同信息场景，验证方案鲁棒性",
    )
    bottom_box = add_box(slide, 3.72, 6.54, 5.12, 0.38, fill=None, line="gray", width=0.9,
                         dash=MSO_LINE_DASH_STYLE.DASH)
    set_shape_text(bottom_box, "形成“指标提出—设备扩展—运行验证”闭环", 15.5, "black", False,
                   PP_ALIGN.CENTER, FONT_BODY, margin=0.02)

    add_box(slide, 9.58, 2.12, 3.3, 4.95, fill=None, line="blue_dash", width=1.6,
            dash=MSO_LINE_DASH_STYLE.DASH)
    add_textbox(slide, 9.88, 2.42, 1.7, 0.28, "实验设计", 21, "blue", True,
                font_name=FONT_HEAD, margin=0.0, valign=MSO_ANCHOR.TOP)
    add_info_section(slide, 9.92, 2.88, "案例场景", "德国社区、松山湖园区")
    add_info_section(slide, 9.92, 3.86, "数据输入", "8760 h 负荷与气象；\n典型日/预测输入")
    add_info_section(slide, 9.92, 4.9, "方法对比", "EQD、Std；\n有/无卡诺电池")
    add_info_section(slide, 9.92, 5.95, "验证重点", "配置差异、运行可行性\n与鲁棒性")
    return slide


def create_standalone_deck():
    global SLIDE_W, SLIDE_H
    template = Presentation(TEMPLATE_PPTX)
    prs = Presentation()
    prs.slide_width = template.slide_width
    prs.slide_height = template.slide_height
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    build_slide_one(prs)
    build_slide_two(prs)
    prs.save(OUTPUT_PPTX)


if __name__ == "__main__":
    create_standalone_deck()
